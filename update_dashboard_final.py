#!/usr/bin/env python3
"""
Enhanced dashboard updater with cash flow forecasting
Preserves guarantee data while adding cash position forecasts
"""
import csv
import re
import os
import glob
import json
from datetime import datetime
from collections import defaultdict
import warnings

# Suppress Pillow and other deprecation warnings that trigger PowerShell errors
warnings.filterwarnings("ignore", category=DeprecationWarning)

# PPTX Dashboard imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_num(v):
    if not v or v.strip() in ('#N/A', '#DIV/0!', '#REF!', ''):
        return 0.0
    try:
        return float(v.strip().replace(',', '').replace('"', ''))
    except ValueError:
        return 0.0

def to_ym(date_str):
    if not date_str or not date_str.strip():
        return None
    return date_str.strip()[:7]

def find_col_by_label(row, label):
    label_lower = label.lower()
    for i, cell in enumerate(row):
        if label_lower in cell.lower().strip():
            if i + 1 < len(row):
                return i + 1
    return None

def parse_collateral_csv(csv_path):
    """Parse the collateral calculator CSV"""
    print(f"  Reading: Collateral calculator...")
    with open(csv_path, newline='', encoding='utf-8-sig') as f:
        rows = list(csv.reader(f))

    fx_rate = 11.0
    utilized_today = 0.0
    cash_collateral = {}

    for ri in range(min(3, len(rows))):
        col = find_col_by_label(rows[ri], 'EUR/NOK')
        if col is not None:
            v = parse_num(rows[ri][col])
            if v > 0:
                fx_rate = v
                break

    for ri in range(min(3, len(rows))):
        col = find_col_by_label(rows[ri], 'Utilized today')
        if col is not None:
            v = parse_num(rows[ri][col])
            if v > 0:
                utilized_today = v
                break

    as_of_date = ''
    if len(rows) > 1 and rows[1]:
        candidate = rows[1][0].strip()
        if re.match(r'^\d{4}-\d{2}-\d{2}$', candidate):
            as_of_date = candidate

    if not as_of_date:
        m = re.search(r'(\d{4}-\d{2}-\d{2})', os.path.basename(csv_path))
        if m:
            as_of_date = m.group(1)

    if not as_of_date:
        from datetime import date
        as_of_date = date.today().strftime('%Y-%m-%d')

    el_volume_date = ''
    el_price_date = ''
    for ri in range(min(5, len(rows))):
        row = rows[ri]
        for ci, cell in enumerate(row):
            if 'el. volume' in cell.lower():
                if ci + 1 < len(row):
                    el_volume_date = row[ci+1].strip()
            if 'el. price' in cell.lower():
                if ci + 1 < len(row):
                    el_price_date = row[ci+1].strip()

    # Extract cash collateral amounts (summed in kNOK)
    for row in rows:
        if len(row) >= 6 and 'cash collateral' in str(row[1]).lower():
            period = to_ym(row[0])
            amount_nok = parse_num(row[5])
            if period:
                cash_collateral[period] = cash_collateral.get(period, 0.0) + amount_nok

    limits = {}
    for row in rows:
        if len(row) > 11:
            period_val = row[7].strip()
            gtype = row[8].strip().lower()
            # DNB facility limit applies to DNB NOK and DNB EUR guarantees (excluding Atradius)
            if period_val and (gtype.startswith('dnb nok') or gtype.startswith('dnb eur')):
                period = to_ym(period_val)
                limit_nok = parse_num(row[11])
                if period:
                    limits[period] = limits.get(period, 0.0) + limit_nok

    actual_marex = {}
    for row in rows:
        if len(row) > 17:
            period_val = row[13].strip()
            gtype = row[14].strip().lower()
            beneficiary = row[15].strip().lower()
            if period_val and gtype == 'dnb eur' and beneficiary == 'marex':
                period = to_ym(period_val)
                amount = parse_num(row[17])
                if period:
                    actual_marex[period] = amount

    forecast = {}
    for row in rows[3:]:
        if len(row) > 5:
            period_val = row[0].strip()
            gtype = row[1].strip()
            beneficiary = row[2].strip()
            amt_k = parse_num(row[4])
            amt_nok = parse_num(row[5])
            
            if not period_val:
                continue
            period = to_ym(period_val)
            if not period:
                continue
                
            if period not in forecast:
                forecast[period] = {
                    'dnbNok': 0.0,
                    'dnbEurLocal': 0.0,
                    'marexEurLocal': 0.0,
                    'atraLocal': 0.0,
                    'cashSkaNok': 0.0,
                    'cashBnpLocal': 0.0,
                    'cashPvnedLocal': 0.0
                }
                
            gtype_lower = gtype.lower()
            ben_lower = beneficiary.lower()
            if 'libra' in ben_lower:
                continue
            
            if gtype_lower == 'dnb nok':
                forecast[period]['dnbNok'] += amt_nok
            elif gtype_lower == 'dnb eur':
                if ben_lower == 'marex':
                    forecast[period]['marexEurLocal'] += amt_k
                else:
                    forecast[period]['dnbEurLocal'] += amt_k
            elif gtype_lower == 'atradius eur':
                forecast[period]['atraLocal'] += amt_k
            elif gtype_lower == 'cash collateral':
                if ben_lower == 'bnp sepa':
                    forecast[period]['cashBnpLocal'] += amt_k
                elif ben_lower == 'pvned':
                    forecast[period]['cashPvnedLocal'] += amt_k
                elif ben_lower == 'skagerak':
                    forecast[period]['cashSkaNok'] += amt_nok

    # Parse Actuals and merge/overwrite forecast values where actuals exist
    actuals = {}
    for row in rows[3:]:
        if len(row) > 18:
            period_val = row[13].strip()
            gtype = row[14].strip()
            beneficiary = row[15].strip()
            amt_k = parse_num(row[17])
            amt_nok = parse_num(row[18])
            
            if not period_val:
                continue
            period = to_ym(period_val)
            if not period:
                continue
                
            if period not in actuals:
                actuals[period] = {
                    'dnbNok': 0.0,
                    'dnbEurLocal': 0.0,
                    'marexEurLocal': 0.0,
                    'atraLocal': 0.0,
                    'cashSkaNok': 0.0,
                    'cashBnpLocal': 0.0,
                    'cashPvnedLocal': 0.0
                }
                
            gtype_lower = gtype.lower()
            ben_lower = beneficiary.lower()
            if 'libra' in ben_lower:
                continue
            
            if gtype_lower == 'dnb nok':
                actuals[period]['dnbNok'] += amt_nok
            elif gtype_lower == 'dnb eur':
                if ben_lower == 'marex':
                    actuals[period]['marexEurLocal'] += amt_k
                else:
                    actuals[period]['dnbEurLocal'] += amt_k
            elif gtype_lower == 'atradius eur':
                actuals[period]['atraLocal'] += amt_k
            elif gtype_lower == 'cash collateral':
                if ben_lower == 'bnp sepa':
                    actuals[period]['cashBnpLocal'] += amt_k
                elif ben_lower == 'pvned':
                    actuals[period]['cashPvnedLocal'] += amt_k
                elif ben_lower == 'skagerak':
                    actuals[period]['cashSkaNok'] += amt_nok



    # Parse and overwrite actual cash collateral in cash_collateral dict
    actual_cash_collateral = {}
    for row in rows:
        if len(row) > 18 and 'cash collateral' in str(row[14]).lower():
            period = to_ym(row[13])
            amount_nok = parse_num(row[18])
            if period:
                actual_cash_collateral[period] = actual_cash_collateral.get(period, 0.0) + amount_nok
                
    for period, amt in actual_cash_collateral.items():
        if amt > 0:
            cash_collateral[period] = amt

    # Parse Today DNB table (columns 39 to 44)
    today_dnb = []
    for row in rows[3:]:
        if len(row) > 44:
            gtype = row[39].strip()
            beneficiary = row[40].strip()
            currency = row[41].strip()
            amt_nok = row[42].strip()
            amt_keur = row[43].strip()
            expiry = row[44].strip()
            
            if (gtype or beneficiary or currency or amt_nok or amt_keur or expiry) and gtype.lower() != 'dnb guarantee':
                today_dnb.append({
                    'gtype': gtype,
                    'beneficiary': beneficiary,
                    'currency': currency,
                    'amt_nok': amt_nok,
                    'amt_keur': amt_keur,
                    'expiry': expiry
                })

    # Dynamically build history dicts from parsed actuals
    history_dnb_nok_only = {}
    history_dnb_eur_local = {}
    history_dnb_nok_base = {}
    history_marex_local = {}

    for period, act in actuals.items():
        as_of_ym = as_of_date[:7] if as_of_date else '2026-08'
        if period > as_of_ym:
            continue
            
        total_dnb_eur = act['dnbEurLocal'] + act['marexEurLocal']
        if total_dnb_eur > 0:
            history_dnb_eur_local[period] = total_dnb_eur
            history_dnb_nok_base[period] = act['dnbNok']
        else:
            if act['dnbNok'] > 0:
                history_dnb_nok_only[period] = act['dnbNok']
                
        if act['marexEurLocal'] > 0:
            history_marex_local[period] = act['marexEurLocal']

    return {
        'fx_rate': fx_rate,
        'utilized_today': utilized_today,
        'as_of_date': as_of_date,
        'cash_collateral': cash_collateral,
        'limits': limits,
        'actual_marex': actual_marex,
        'forecast': forecast,
        'today_dnb': today_dnb,
        'history_dnb_nok_only': history_dnb_nok_only,
        'history_dnb_eur_local': history_dnb_eur_local,
        'history_dnb_nok_base': history_dnb_nok_base,
        'history_marex_local': history_marex_local,
        'actuals': actuals,
        'el_volume_date': el_volume_date,
        'el_price_date': el_price_date
    }

def parse_treasury_csv(csv_path):
    """Parse CFF Treasury CSV (weekly cash forecasts)"""
    print(f"  Reading: CFF Treasury (weekly)...")
    with open(csv_path, newline='', encoding='utf-8-sig') as f:
        rows = list(csv.reader(f))

    treasury_data = {}

    # Skip header rows (rows 0-2)
    for row in rows[2:]:
        if len(row) >= 3 and row[1].strip():
            try:
                week_num = int(row[1].strip())
                import datetime
                # Calculate correct Monday of that week in 2026 using ISO week formula
                monday_date = datetime.datetime.strptime(f"2026-W{week_num:02d}-1", "%G-W%V-%u").date()
                date_str = monday_date.strftime("%Y-%m-%d")
                cash_pos = parse_num(row[2]) * 1000  # Current forecast column (converted from mNOK to kNOK)
                treasury_data[date_str] = cash_pos
            except:
                pass

    return treasury_data

def parse_fpa_csv(csv_path):
    """Parse CFF FP&A CSV (monthly cash forecasts)"""
    print(f"  Reading: CFF FP&A (monthly)...")
    with open(csv_path, newline='', encoding='utf-8-sig') as f:
        rows = list(csv.reader(f))

    fpa_data = {}

    # Skip header rows, start from row 5
    for row in rows[5:]:
        if len(row) >= 2 and row[0].strip():
            try:
                date_str = row[0].strip()
                # Parse dates like "01/04/2025" to "2025-04-01"
                if '/' in date_str:
                    parts = date_str.split('/')
                    if len(parts) == 3:
                        day, month, year = parts
                        date_str = f"{year}-{month}-{day}"

                if re.match(r'^\d{4}-\d{2}-\d{2}$', date_str):
                    fpa_data[date_str] = {
                        'BU26': parse_num(row[5]) if len(row) >= 6 else parse_num(row[1]),
                        'FC1.5': parse_num(row[6]) if len(row) >= 7 else 0.0,
                        'Actual': parse_num(row[7]) if len(row) >= 8 else 0.0,
                        'WCFF': parse_num(row[8]) if len(row) >= 9 else 0.0
                    }
            except:
                pass

    return fpa_data

def bootstrap_cash_flow(treasury, fpa, current_month=None):
    """Bootstrap weekly treasury data to monthly and bridge with FP&A"""
    monthly_cash = {}
    cash_sources = {}

    if not current_month:
        current_month = "2026-08"

    # Determine current 3 months (current, current + 1, current + 2)
    try:
        import datetime
        current_date = datetime.datetime.strptime(current_month + "-01", "%Y-%m-%d").date()
        current_3_months = []
        for i in range(3):
            m = current_date.month + i
            y = current_date.year
            if m > 12:
                m -= 12
                y += 1
            current_3_months.append(f"{y}-{m:02d}")
    except:
        current_3_months = ["2026-08", "2026-09", "2026-10"]

    # Process FP&A data to populate monthly cash positions and sources based on rules
    for date_str, values in fpa.items():
        month_str = date_str[:7]
        
        # Rule (1): Earlier than current month -> use Actual from CFF FP&A
        if month_str < current_month:
            val = values.get('Actual', 0.0)
            if val > 0:
                monthly_cash[month_str] = val
                cash_sources[month_str] = "Actual"
                
        # Rule (2): Current 3 months -> use WCFF from CFF FP&A
        elif month_str in current_3_months:
            val = values.get('WCFF', 0.0)
            if val > 0:
                monthly_cash[month_str] = val
                cash_sources[month_str] = "WCFF"
                
        # Rule (3): Future months in 2026 -> use BU26 from CFF FP&A (rename to Budget26)
        elif month_str.startswith('2026-'):
            val = values.get('BU26', 0.0)
            if val > 0:
                monthly_cash[month_str] = val
                cash_sources[month_str] = "Budget26"
                
        # Rule (4): Future months in 2027 and later -> use FC 1.5 from CFF FP&A
        elif month_str >= '2027-01':
            val = values.get('FC1.5', 0.0)
            if val > 0:
                monthly_cash[month_str] = val
                cash_sources[month_str] = "FC 1.5"

    return monthly_cash, cash_sources

def format_forecast_js(data):
    """Format data dict as JavaScript object literal"""
    if not data:
        return '{}'
    lines = []
    for key in sorted(data.keys()):
        lines.append(f'  "{key}": {int(data[key])},')
    if lines:
        lines[-1] = lines[-1].rstrip(',')  # Remove trailing comma from last item
    return '{\n' + '\n'.join(lines) + '\n}'

def format_sources_js(data):
    """Format data dict with string values as JavaScript object literal"""
    if not data:
        return '{}'
    lines = []
    for key in sorted(data.keys()):
        lines.append(f'  "{key}": "{data[key]}",')
    if lines:
        lines[-1] = lines[-1].rstrip(',')
    return '{\n' + '\n'.join(lines) + '\n}'

def format_forecast_block(forecast_dict):
    """Format forecast dict of dicts as JavaScript object literal"""
    if not forecast_dict:
        return '{}'
    lines = []
    for key in sorted(forecast_dict.keys()):
        val = forecast_dict[key]
        parts = []
        for k in ['dnbNok', 'dnbEurLocal', 'marexEurLocal', 'atraLocal', 'cashSkaNok', 'cashBnpLocal', 'cashPvnedLocal']:
            v = val.get(k, 0.0)
            parts.append(f'{k}: {int(round(v))}')
        line = f'  "{key}": {{ ' + ', '.join(parts) + ' },'
        lines.append(line)
    if lines:
        lines[-1] = lines[-1].rstrip(',')
    return '{\n' + '\n'.join(lines) + '\n}'

def extract_existing_data(html_content):
    """Extract existing FORECAST, HISTORY data from HTML"""
    data = {}

    # Extract the DATA block
    match = re.search(r'// DATA_START\n(.*?)\n// DATA_END', html_content, re.DOTALL)
    if match:
        block = match.group(1)

        # Extract FORECAST object
        forecast_match = re.search(r'const FORECAST = (\{.*?\});', block, re.DOTALL)
        if forecast_match:
            try:
                # Parse as JavaScript object (simplified)
                data['FORECAST'] = forecast_match.group(1)
            except:
                pass

        # Extract HISTORY data
        for var in ['HISTORY_DNB_NOK_ONLY', 'HISTORY_DNB_EUR_LOCAL', 'HISTORY_DNB_NOK_BASE', 'HISTORY_MAREX_LOCAL']:
            var_match = re.search(rf'const {var} = (\{{.*?\}});', block, re.DOTALL)
            if var_match:
                try:
                    data[var] = var_match.group(1)
                except:
                    pass

    return data

def parse_js_obj(js_str):
    # Quote keys
    json_str = re.sub(r'([a-zA-Z0-9_]+):', r'"\1":', js_str)
    # Remove trailing commas
    json_str = re.sub(r',\s*\}', r'}', json_str)
    json_str = re.sub(r',\s*\]', r']', json_str)
    return json.loads(json_str)

def update_html(html_content, collateral_data, monthly_cash, cash_collateral, cash_sources):
    """Update HTML preserving guarantee data and adding cash flow data"""

    # Extract existing guarantee data
    existing = extract_existing_data(html_content)

    # Build new data block with all data
    forecast_dict = collateral_data['forecast']
    forecast_obj = format_forecast_block(forecast_dict)
    history_dnb_nok = format_forecast_js(collateral_data['history_dnb_nok_only'])
    history_dnb_eur = format_forecast_js(collateral_data['history_dnb_eur_local'])
    history_dnb_nok_base = format_forecast_js(collateral_data['history_dnb_nok_base'])
    history_marex = format_forecast_js(collateral_data['history_marex_local'])

    # Calculate cash after collateral
    fx_rate = collateral_data['fx_rate']
    limits = collateral_data['limits']

    as_of_period = collateral_data['as_of_date'][:7] if collateral_data['as_of_date'] else ""
    cash_coll_today = cash_collateral.get(as_of_period, 0.0) if as_of_period else 0.0

    cash_after_collateral = {}
    for period, cash_pos in monthly_cash.items():
        fd = forecast_dict.get(period, {})
        dnb_nok = fd.get('dnbNok', 0.0)
        dnb_eur = fd.get('dnbEurLocal', 0.0) * fx_rate
        marex = fd.get('marexEurLocal', 0.0) * fx_rate
        atra = fd.get('atraLocal', 0.0) * fx_rate

        utilization = dnb_nok + dnb_eur + marex + atra
        limit = limits.get(period, 200000.0)

        excess = max(0.0, utilization - limit)
        coll_increase = cash_collateral.get(period, 0.0) - cash_coll_today

        cash_after_collateral[period] = cash_pos - excess - coll_increase

    # Build complete data block
    today_dnb_json = json.dumps(collateral_data['today_dnb'])
    actuals_obj = format_forecast_block(collateral_data['actuals'])
    data_block = f"""// DATA_START
const BASE_FX = {collateral_data['fx_rate']};
const FORECAST = {forecast_obj};
const ACTUALS = {actuals_obj};
const HISTORY_DNB_NOK_ONLY = {history_dnb_nok};
const HISTORY_DNB_EUR_LOCAL = {history_dnb_eur};
const HISTORY_DNB_NOK_BASE = {history_dnb_nok_base};
const HISTORY_MAREX_LOCAL = {history_marex};
const MONTHLY_CASH_FORECAST = {format_forecast_js(monthly_cash)};
const CASH_COLLATERAL = {format_forecast_js(cash_collateral)};
const CASH_AFTER_COLLATERAL = {format_forecast_js(cash_after_collateral)};
const CREDIT_LIMIT = {format_forecast_js(collateral_data['limits'])};
const CASH_SOURCES = {format_sources_js(cash_sources)};
const EL_VOLUME_DATE = "{collateral_data.get('el_volume_date', '')}";
const EL_PRICE_DATE = "{collateral_data.get('el_price_date', '')}";
const TODAY_DNB = {today_dnb_json};
// DATA_END"""

    # Replace the entire DATA block
    html_content = re.sub(
        r'// DATA_START.*?// DATA_END',
        data_block,
        html_content,
        flags=re.DOTALL
    )

    # Update KPI
    utilized_str = '{:,}'.format(round(collateral_data['utilized_today']))
    html_content = re.sub(
        r'(id="kpi-utilized"[^>]*>)[^<]*',
        lambda m: m.group(1) + utilized_str,
        html_content
    )

    # Update as-of date
    as_of = collateral_data['as_of_date']
    if as_of:
        html_content = re.sub(
            r'As of:?\s*\d{4}-\d{2}-\d{2}',
            'As of: ' + as_of,
            html_content
        )
        
        # Update actuals month label (e.g. "Jun actuals" to "Aug actuals")
        try:
            as_of_dt = datetime.strptime(as_of, "%Y-%m-%d")
            month_name = as_of_dt.strftime("%b")
            html_content = re.sub(
                r'class="kpi-sub">kNOK\s*\([A-Za-z]+\s*actuals?\)',
                f'class="kpi-sub">kNOK ({month_name} actuals)',
                html_content
            )
        except Exception as e:
            print(f"Warning: could not parse as_of date '{as_of}' for KPI sub-label: {e}")

    return html_content

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_slide_header(slide, title_text, category_text):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 208, 231)  # Tibber Blue
    p.font.name = 'Segoe UI'

    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(8), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(22, 25, 29)  # Soft Black
    p2.font.name = 'Segoe UI'

def add_report_date_footer(slide, as_of_date):
    txBox3 = slide.shapes.add_textbox(Inches(9.5), Inches(6.9), Inches(3.2), Inches(0.4))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_top = tf3.margin_bottom = tf3.margin_right = 0
    p3 = tf3.paragraphs[0]
    p3.text = f"Report Date: {as_of_date}"
    p3.alignment = PP_ALIGN.RIGHT
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(104, 108, 115)  # Dark Grey
    p3.font.name = 'Segoe UI'

def add_card(slide, left, top, width, height, title, value_str, subtext=None, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.color.rgb = RGBColor(212, 217, 224)  # Medium Grey
        card.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title if "kNOK" in title else title.upper()
    p_title.font.size = Pt(10)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(104, 108, 115)  # Dark Grey
    p_title.font.name = 'Segoe UI'
    p_title.space_after = Pt(8)
    
    p_val = tf.add_paragraph()
    p_val.text = value_str
    p_val.font.size = Pt(28)
    p_val.font.bold = True
    p_val.font.color.rgb = RGBColor(22, 25, 29)  # Soft Black
    p_val.font.name = 'Segoe UI'
    
    if subtext:
        p_val.space_after = Pt(4)
        p_sub = tf.add_paragraph()
        p_sub.text = subtext
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = RGBColor(104, 108, 115)  # Dark Grey
        p_sub.font.name = 'Segoe UI'

def make_white_logo(original_path, white_path):
    """Invert colors of transparent PNG to create white version of Tibber logo"""
    try:
        from PIL import Image
        img = Image.open(original_path).convert("RGBA")
        data = img.getdata()
        new_data = []
        for item in data:
            if item[3] > 0:
                # Replace with white (255, 255, 255) preserving original alpha (transparency)
                new_data.append((255, 255, 255, item[3]))
            else:
                new_data.append(item)
        img.putdata(new_data)
        img.save(white_path)
    except Exception as e:
        print(f"Warning: failed to make white logo: {e}")

def generate_pptx_dashboard(collateral_data, monthly_cash, cash_collateral, pptx_path):
    """Generate modern, professionally styled slide deck for the collateral dashboard"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_path = os.path.join(script_dir, "tibber_logo.png")
    white_logo_path = os.path.join(script_dir, "tibber_logo_white.png")
    
    has_logo = os.path.exists(logo_path)
    if has_logo:
        make_white_logo(logo_path, white_logo_path)
        
    # Extract data from HTML to replicate calculations and tables
    forecast_dict = {}
    history_dnb_nok_only = {}
    history_dnb_eur_local = {}
    history_dnb_nok_base = {}
    history_marex_local = {}
    html_path = pptx_path.replace(".pptx", ".html")
    if os.path.exists(html_path):
        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            existing = extract_existing_data(html_content)
            forecast_dict = parse_js_obj(existing.get('FORECAST', '{}'))
            history_dnb_nok_only = parse_js_obj(existing.get('HISTORY_DNB_NOK_ONLY', '{}'))
            history_dnb_eur_local = parse_js_obj(existing.get('HISTORY_DNB_EUR_LOCAL', '{}'))
            history_dnb_nok_base = parse_js_obj(existing.get('HISTORY_DNB_NOK_BASE', '{}'))
            history_marex_local = parse_js_obj(existing.get('HISTORY_MAREX_LOCAL', '{}'))
        except Exception as e:
            print(f"Warning: PPTX generator failed to read existing HTML: {e}")

    fx = collateral_data['fx_rate']
    limits = collateral_data['limits']
    periods = sorted(forecast_dict.keys())
    
    dnb_nok_arr = []
    dnb_eur_arr = []
    marex_arr = []
    atra_arr = []
    cash_ska_arr = []
    cash_bnp_arr = []
    cash_pvned_arr = []
    dnb_total_arr = []
    cash_total_arr = []
    grand_total_arr = []
    credit_limit_arr = []
    
    for p in periods:
        fd = forecast_dict.get(p, {})
        dnb_nok = fd.get('dnbNok', 0.0)
        dnb_eur = fd.get('dnbEurLocal', 0.0) * fx
        marex = fd.get('marexEurLocal', 0.0) * fx
        atra = fd.get('atraLocal', 0.0) * fx
        
        cash_ska = fd.get('cashSkaNok', 0.0)
        cash_bnp = fd.get('cashBnpLocal', 0.0) * fx
        cash_pvned = fd.get('cashPvnedLocal', 0.0) * fx
        
        dnb_total = dnb_nok + dnb_eur + marex
        cash_total = cash_ska + cash_bnp + cash_pvned
        
        dnb_nok_arr.append(dnb_nok)
        dnb_eur_arr.append(dnb_eur)
        marex_arr.append(marex)
        atra_arr.append(atra)
        cash_ska_arr.append(cash_ska)
        cash_bnp_arr.append(cash_bnp)
        cash_pvned_arr.append(cash_pvned)
        dnb_total_arr.append(dnb_total)
        cash_total_arr.append(cash_total)
        grand_total_arr.append(dnb_total + atra + cash_total)
        
        limit = limits.get(p, 200000.0)
        credit_limit_arr.append(limit)

    as_of_date = collateral_data.get('as_of_date', 'N/A')
    as_of_period = as_of_date[:7] if as_of_date else ""
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Tibber Dark theme)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1, RGBColor(22, 25, 29)) # Soft Black
    
    if has_logo and os.path.exists(white_logo_path):
        slide1.shapes.add_picture(white_logo_path, Inches(1.0), Inches(0.8), width=Inches(1.8))
        
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_sub1 = tf.paragraphs[0]
    p_sub1.text = "TREASURY & FINANCE REPORT"
    p_sub1.font.size = Pt(12)
    p_sub1.font.bold = True
    p_sub1.font.color.rgb = RGBColor(30, 208, 231)  # Tibber Blue
    p_sub1.font.name = 'Segoe UI'
    p_sub1.space_after = Pt(12)
    
    p_title = tf.add_paragraph()
    p_title.text = "Guarantee & Collateral Dashboard"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    p_title.font.name = 'Segoe UI'
    p_title.space_after = Pt(8)
    
    p_sub2 = tf.add_paragraph()
    p_sub2.text = "Automated forecast and risk limits overview"
    p_sub2.font.size = Pt(16)
    p_sub2.font.color.rgb = RGBColor(212, 217, 224)  # Medium Grey
    p_sub2.font.name = 'Segoe UI'
    p_sub2.space_after = Pt(40)
    
    p_date = tf.add_paragraph()
    p_date.text = f"As of Date: {as_of_date}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    p_date.font.size = Pt(12)
    p_date.font.color.rgb = RGBColor(104, 108, 115)  # Dark Grey
    p_date.font.name = 'Segoe UI'
    
    # ----------------------------------------------------
    # SLIDE 2: KPI Metrics Dashboard (Tibber Light theme)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, RGBColor(251, 251, 251)) # Off White
    add_slide_header(slide2, "Key Performance Indicators", "Overview")
    add_report_date_footer(slide2, as_of_date)
    
    if has_logo:
        slide2.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), width=Inches(1.5))
        
    limit_val = limits.get(as_of_period, 200000.0) if as_of_period else 200000.0
    cash_coll_val = cash_collateral.get(as_of_period, 0.0) if as_of_period else 0.0
    
    # Compute Breach KPI
    breach_periods = []
    for i, p in enumerate(periods):
        excess = dnb_total_arr[i] - credit_limit_arr[i]
        if excess > 0:
            breach_periods.append((p, excess))
            
    if breach_periods:
        first_breach = breach_periods[0]
        breach_val_str = f"+{first_breach[1]:,.0f} kNOK"
        breach_sub = f"first breach · {first_breach[0]}"
        breach_border = RGBColor(255, 90, 90) # Red Border for breach
    else:
        breach_val_str = "Within Limit"
        breach_sub = "no breach in forecast"
        breach_border = RGBColor(30, 208, 231) # Blue Border
        
    # Compute Peak DNB KPI
    if dnb_total_arr:
        peak_dnb_val = max(dnb_total_arr)
        peak_dnb_idx = dnb_total_arr.index(peak_dnb_val)
        peak_dnb_period = periods[peak_dnb_idx]
    else:
        peak_dnb_val = 0
        peak_dnb_period = "N/A"
        
    # Compute Peak Grand Total KPI
    if grand_total_arr:
        peak_grand_val = max(grand_total_arr)
        peak_grand_idx = grand_total_arr.index(peak_grand_val)
        peak_grand_period = periods[peak_grand_idx]
    else:
        peak_grand_val = 0
        peak_grand_period = "N/A"
        
    # Row 1 Cards
    add_card(slide2, Inches(0.8), Inches(1.5), Inches(3.6), Inches(1.8), 
             "Today DNB Guarantee Utilized", f"{collateral_data['utilized_today']:,.0f} kNOK", "kNOK (Jun actuals)", RGBColor(30, 208, 231))
             
    add_card(slide2, Inches(4.8), Inches(1.5), Inches(3.6), Inches(1.8), 
             "DNB Guarantee Line", f"{limit_val:,.0f} kNOK", f"Facility credit limit for {as_of_period or 'current month'}")
             
    add_card(slide2, Inches(8.8), Inches(1.5), Inches(3.6), Inches(1.8), 
             "DNB Guarantee BREACH (kNOK)", breach_val_str, breach_sub, breach_border)
             
    # Row 2 Cards
    add_card(slide2, Inches(0.8), Inches(3.5), Inches(3.6), Inches(1.8), 
             "PEAK DNB Guarantee FORECAST (kNOK)", f"{peak_dnb_val:,.0f} kNOK", f"Peak period: {peak_dnb_period}")
             
    add_card(slide2, Inches(4.8), Inches(3.5), Inches(3.6), Inches(1.8), 
             "PEAK ALL guarantee/collaterals TOTAL (kNOK)", f"{peak_grand_val:,.0f} kNOK", f"DNB + Atradius + Cash · Peak: {peak_grand_period}")
             
    add_card(slide2, Inches(8.8), Inches(3.5), Inches(3.6), Inches(1.8), 
             "EUR/NOK FX Rate", f"{fx:.4f}", "Used for currency conversions")
             
    # Bottom callout card
    bottom_card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.6), Inches(1.2))
    bottom_card.fill.solid()
    bottom_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bottom_card.line.color.rgb = RGBColor(212, 217, 224)
    
    tb_bottom = slide2.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.2), Inches(1.0))
    tf_bottom = tb_bottom.text_frame
    tf_bottom.word_wrap = True
    tf_bottom.margin_left = tf_bottom.margin_right = tf_bottom.margin_top = tf_bottom.margin_bottom = 0
    
    p_b_title = tf_bottom.paragraphs[0]
    p_b_title.text = f"SUMMARY: Cash Collateral holding is {cash_coll_val:,.0f} kNOK. Available Credit line headroom: {(limit_val - collateral_data['utilized_today']):,.0f} kNOK."
    p_b_title.font.size = Pt(11)
    p_b_title.font.bold = True
    p_b_title.font.color.rgb = RGBColor(22, 25, 29)
    p_b_title.font.name = 'Segoe UI'
    
    p_b_desc = tf_bottom.add_paragraph()
    p_b_desc.text = "All cash forecasts are bridged between weekly CFF Treasury and monthly CFF FP&A data sources. Covenant limit is monitored at 125,000 kNOK."
    p_b_desc.font.size = Pt(10.5)
    p_b_desc.font.color.rgb = RGBColor(104, 108, 115)
    p_b_desc.font.name = 'Segoe UI'
    p_b_desc.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 3: Cash Calculation Detail Table (Panel 2B)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, RGBColor(251, 251, 251))
    add_slide_header(slide3, "Calculation of Cash After Collateral", "Cash Flow Impact")
    add_report_date_footer(slide3, as_of_date)
    
    if has_logo:
        slide3.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), width=Inches(1.5))
        
    # Get sorted months that have cash forecast from current months
    as_of_period = collateral_data.get('as_of_date', '')[:7] if collateral_data.get('as_of_date') else "2026-06"
    sorted_months = sorted([p for p in periods if p in monthly_cash and p >= as_of_period])
    if not sorted_months:
        sorted_months = sorted([p for p in periods if p in monthly_cash])[:8]
    
    if sorted_months:
        rows = 5 # Metric name + 4 value rows
        cols = len(sorted_months) + 1
        
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.7)
        height = Inches(3.6)
        
        table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Column widths
        col0_w = 1.9 if cols > 10 else 2.7
        table.columns[0].width = Inches(col0_w)
        col_w = (11.7 - col0_w) / len(sorted_months)
        for c in range(1, cols):
            table.columns[c].width = Inches(col_w)
            
        # Headers
        cell = table.cell(0, 0)
        cell.text = "Metric (kNOK)"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(22, 25, 29)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8.5) if cols > 10 else Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Segoe UI'
        
        for c, p_ym in enumerate(sorted_months):
            cell = table.cell(0, c + 1)
            cell.text = p_ym
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(22, 25, 29)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(8.5) if cols > 10 else Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Segoe UI'
            
        # Rows
        row_labels = [
            "Cash Position",
            "Guarantee Overdraft (Breach)",
            "Cash Collateral Change",
            "Cash After Collateral"
        ]
        
        cash_coll_today = cash_collateral.get(as_of_period, 0.0) if as_of_period else 0.0
        
        for r_idx, label in enumerate(row_labels):
            row_num = r_idx + 1
            cell = table.cell(row_num, 0)
            cell.text = label
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(243, 246, 249) if row_num % 2 == 1 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8) if cols > 10 else Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(22, 25, 29)
            p.font.name = 'Segoe UI'
            if r_idx == 3: # Highlight Cash After Collateral row label
                p.font.color.rgb = RGBColor(16, 185, 129)
                
            for c_idx, p_ym in enumerate(sorted_months):
                cell = table.cell(row_num, c_idx + 1)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(243, 246, 249) if row_num % 2 == 1 else RGBColor(255, 255, 255)
                
                # Fetch/calculate values
                cash_pos = monthly_cash.get(p_ym, 0.0)
                fd = forecast_dict.get(p_ym, {})
                dnb_nok = fd.get('dnbNok', 0.0)
                dnb_eur = fd.get('dnbEurLocal', 0.0) * fx
                marex = fd.get('marexEurLocal', 0.0) * fx
                atra = fd.get('atraLocal', 0.0) * fx
                utilization = dnb_nok + dnb_eur + marex + atra
                limit = limits.get(p_ym, 200000.0)
                excess = max(0.0, utilization - limit)
                
                cash_coll = cash_collateral.get(p_ym, 0.0)
                coll_change = cash_coll - cash_coll_today
                
                # Formula matching HTML: cash_after = cash_pos - excess - coll_change
                cash_after = cash_pos - excess - coll_change
                
                val_str = ""
                cell_color = RGBColor(104, 108, 115)
                is_bold = False
                
                if r_idx == 0:
                    val_str = f"{cash_pos:,.0f}"
                elif r_idx == 1:
                    val_str = f"+{excess:,.0f}" if excess > 0 else "-"
                    if excess > 0:
                        cell_color = RGBColor(255, 90, 90)
                        is_bold = True
                elif r_idx == 2:
                    val_str = f"+{coll_change:,.0f}" if coll_change > 0 else (f"{coll_change:,.0f}" if coll_change < 0 else "0")
                    if coll_change > 0:
                        cell_color = RGBColor(255, 90, 90)
                    elif coll_change < 0:
                        cell_color = RGBColor(16, 185, 129)
                elif r_idx == 3:
                    val_str = f"{cash_after:,.0f}"
                    cell_color = RGBColor(16, 185, 129)
                    is_bold = True
                    
                cell.text = val_str
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(8) if cols > 10 else Pt(10)
                p.font.bold = is_bold
                p.font.color.rgb = cell_color
                p.font.name = 'Segoe UI'

    # ----------------------------------------------------
    # SLIDE 4: Forecast vs Actual Deviation (Panel 5)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, RGBColor(251, 251, 251))
    add_slide_header(slide4, "DNB Guarantee Utilization: Forecast vs Actual Deviation", "Risk Management & Variance")
    add_report_date_footer(slide4, as_of_date)
    
    if has_logo:
        slide4.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), width=Inches(1.5))
        
    # Compute history DNB total
    
    def compute_history_dnb(history_dnb_nok_only, history_dnb_eur_local, history_dnb_nok_base, fx):
        result = {}
        for p_ym, v in history_dnb_nok_only.items():
            result[p_ym] = v
        for p_ym, eur_local in history_dnb_eur_local.items():
            result[p_ym] = history_dnb_nok_base.get(p_ym, 0.0) + eur_local * fx
        return result

    hist_dnb = compute_history_dnb(history_dnb_nok_only, history_dnb_eur_local, history_dnb_nok_base, fx)
    overlapping = [p for p in periods if p in hist_dnb]
    
    if overlapping:
        rows = len(overlapping) + 1
        cols = 5
        
        left = Inches(1.5)
        top = Inches(1.8)
        width = Inches(10.3)
        height = Inches(0.45 * rows)
        
        table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        table.columns[0].width = Inches(1.8) # Month
        table.columns[1].width = Inches(2.1) # Forecast
        table.columns[2].width = Inches(2.1) # Actual
        table.columns[3].width = Inches(2.5) # Deviation
        table.columns[4].width = Inches(1.8) # Trend
        
        headers = ["Month", "Forecast (kNOK)", "Actual (kNOK)", "Deviation (kNOK)", "Trend"]
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(22, 25, 29)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Segoe UI'
            
        for r_idx, p_ym in enumerate(overlapping):
            row_num = r_idx + 1
            bg_color = RGBColor(255, 255, 255) if r_idx % 2 == 0 else RGBColor(243, 246, 249)
            
            # Values
            p_idx = periods.index(p_ym)
            fc_val = dnb_total_arr[p_idx]
            hist_val = hist_dnb[p_ym]
            dev = fc_val - hist_val
            dev_pct = (dev / hist_val * 100) if hist_val else 0.0
            
            trend = "-"
            if r_idx > 0:
                prev_p = overlapping[r_idx - 1]
                prev_fc = dnb_total_arr[periods.index(prev_p)]
                prev_hist = hist_dnb[prev_p]
                prev_dev = prev_fc - prev_hist
                if dev > prev_dev + 500:
                    trend = "Rising"
                elif dev < prev_dev - 500:
                    trend = "Falling"
                else:
                    trend = "Stable"
                    
            dev_str = f"+{dev:,.0f} (+{dev_pct:.1f}%)" if dev > 0 else (f"{dev:,.0f} ({dev_pct:.1f}%)" if dev < 0 else "0 (0.0%)")
            dev_color = RGBColor(255, 90, 90) if dev > 0 else (RGBColor(30, 208, 231) if dev < 0 else RGBColor(104, 108, 115))
            
            row_data = [
                p_ym,
                f"{fc_val:,.0f}",
                f"{hist_val:,.0f}",
                dev_str,
                trend
            ]
            
            for c_idx, val in enumerate(row_data):
                cell = table.cell(row_num, c_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                p.font.size = Pt(10)
                p.font.name = 'Segoe UI'
                p.font.color.rgb = RGBColor(22, 25, 29) if c_idx != 3 else dev_color
                if c_idx == 0:
                    p.font.bold = True
                if c_idx == 3:
                    p.font.bold = True

    # ----------------------------------------------------
    # SLIDE 5: Supporting Views Breakdown
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, RGBColor(251, 251, 251))
    add_slide_header(slide5, "Guarantees & Cash Detailed Forecast Breakdown", "Supporting Views")
    add_report_date_footer(slide5, as_of_date)
    
    if has_logo:
        slide5.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), width=Inches(1.5))
        
    as_of_period = collateral_data.get('as_of_date', '')[:7] if collateral_data.get('as_of_date') else "2026-06"
    sorted_months = sorted([p for p in periods if p >= as_of_period])
    if not sorted_months:
        sorted_months = periods[:8]
        
    if sorted_months:
        rows = len(sorted_months) + 1
        cols = 7
        
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.7)
        row_h = 0.22 if rows > 10 else 0.45
        height = Inches(row_h * rows)
        
        table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        for r_idx in range(rows):
            table.rows[r_idx].height = Inches(row_h)
        
        # Adjust column widths
        table.columns[0].width = Inches(1.4) # Period
        table.columns[1].width = Inches(1.7) # DNB NOK
        table.columns[2].width = Inches(1.7) # DNB EUR
        table.columns[3].width = Inches(1.7) # Marex
        table.columns[4].width = Inches(1.7) # Atradius
        table.columns[5].width = Inches(1.7) # Cash Collateral
        table.columns[6].width = Inches(1.8) # Grand Total
        
        headers = ["Period", "DNB NOK (kNOK)", "DNB EUR (conv.)", "Marex (conv.)", "Atradius (conv.)", "Cash Collateral", "Grand Total"]
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(22, 25, 29)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.font.size = Pt(8.5) if rows > 10 else Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Segoe UI'
            
        for r_idx, p_ym in enumerate(sorted_months):
            row_num = r_idx + 1
            bg_color = RGBColor(255, 255, 255) if r_idx % 2 == 0 else RGBColor(243, 246, 249)
            
            p_idx = periods.index(p_ym)
            
            # Fetch breakdown details
            dnb_nok = dnb_nok_arr[p_idx]
            dnb_eur = dnb_eur_arr[p_idx]
            marex = marex_arr[p_idx]
            atra = atra_arr[p_idx]
            cash_coll = cash_collateral.get(p_ym, 0.0)
            g_total = grand_total_arr[p_idx]
            
            row_data = [
                p_ym,
                f"{dnb_nok:,.0f}",
                f"{dnb_eur:,.0f}" if dnb_eur > 0 else "-",
                f"{marex:,.0f}" if marex > 0 else "-",
                f"{atra:,.0f}" if atra > 0 else "-",
                f"{cash_coll:,.0f}" if cash_coll > 0 else "-",
                f"{g_total:,.0f}"
            ]
            
            for c_idx, val in enumerate(row_data):
                cell = table.cell(row_num, c_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                p.font.size = Pt(8) if rows > 10 else Pt(9.5)
                p.font.name = 'Segoe UI'
                p.font.color.rgb = RGBColor(22, 25, 29)
                if c_idx == 0:
                    p.font.bold = True
                if c_idx == 6:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(30, 208, 231) # highlight grand total
                    
    # ----------------------------------------------------
    # SLIDE 6: Human Value & Google Antigravity Collaboration (Tibber Light theme)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6, RGBColor(251, 251, 251))
    add_slide_header(slide6, "Human Value & AI Collaboration in Automation", "Governance & Tools")
    add_report_date_footer(slide6, as_of_date)
    
    if has_logo:
        slide6.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), width=Inches(1.5))
        
    # Left Box: Human Value & Responsibility
    left_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_card.line.color.rgb = RGBColor(30, 208, 231)  # Tibber Blue
    left_card.line.width = Pt(1.5)
    
    tb_left = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
    
    p_lh = tf_left.paragraphs[0]
    p_lh.text = "HUMAN VALUE & RESPONSIBILITY"
    p_lh.font.size = Pt(13)
    p_lh.font.bold = True
    p_lh.font.color.rgb = RGBColor(22, 25, 29)
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(14)
    
    bullets_left = [
        ("Verify & Maintain Accuracy", "Ensure data sets are clean and accurate. While AI automates data fetching and slicing, human verification acts as the ultimate quality gate for financial reports."),
        ("Define Linkages & Context", "Establish the logical connections between disparate data sets (e.g., matching weekly treasury forecasts with monthly budget allocations) that AI cannot infer alone."),
        ("Direct Strategic Purpose", "Determine the goal, audience, and narrative direction. AI is a powerful assistant, but only humans know the 'what' and 'why' behind the presentation objectives.")
    ]
    
    for title, desc in bullets_left:
        p_title = tf_left.add_paragraph()
        p_title.text = f"•  {title}"
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(22, 25, 29)
        p_title.font.name = 'Segoe UI'
        p_title.space_before = Pt(8)
        
        p_desc = tf_left.add_paragraph()
        p_desc.text = f"   {desc}"
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = RGBColor(104, 108, 115)
        p_desc.font.name = 'Segoe UI'
        p_desc.space_after = Pt(6)
        
    # Right Box: Why Google Antigravity
    right_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_card.line.color.rgb = RGBColor(212, 217, 224)
    right_card.line.width = Pt(1)
    
    tb_right = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0
    
    p_rh = tf_right.paragraphs[0]
    p_rh.text = "WHY WE UTILIZE GOOGLE ANTIGRAVITY"
    p_rh.font.size = Pt(13)
    p_rh.font.bold = True
    p_rh.font.color.rgb = RGBColor(22, 25, 29)
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(14)
    
    bullets_right = [
        ("Seamless Workspace Integration", "Directly accesses files in Google Drive and downloads/processes data from Google Sheets natively with unmatched reliability."),
        ("Multi-Model Capabilities", "Supports multiple models (Claude, ChatGPT, Google models) in a unified environment, giving flexibility and optimal token usage for automation pipelines."),
        ("Superior Performance & Accuracy", "Demonstrates significantly higher speed, reliability, and precision in parsing and transforming financial data compared to other AI tools.")
    ]
    
    for title, desc in bullets_right:
        p_title = tf_right.add_paragraph()
        p_title.text = f"•  {title}"
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(22, 25, 29)
        p_title.font.name = 'Segoe UI'
        p_title.space_before = Pt(8)
        
        p_desc = tf_right.add_paragraph()
        p_desc.text = f"   {desc}"
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = RGBColor(104, 108, 115)
        p_desc.font.name = 'Segoe UI'
        p_desc.space_after = Pt(6)

    try:
        prs.save(pptx_path)
        print(f"  - PPTX saved successfully: {pptx_path}")
    except PermissionError:
        print(f"\nERROR: Permission denied writing to '{pptx_path}'.")
        print("       Please close the PowerPoint file if it is open in PowerPoint and try again.\n")
        import sys
        sys.exit(1)



def find_input_csvs(script_dir, target_date=None):
    """
    Search for collateral, treasury, and FP&A CSV files across multiple candidate directories
    with case-insensitive matching for full cross-platform / GitHub Actions compatibility.
    """
    candidate_dirs = [
        os.path.join(script_dir, "Data"),
        os.path.join(script_dir, "data"),
        os.path.join(script_dir, "DATA"),
        script_dir,
        os.path.join(script_dir, "Data", "Older"),
        os.path.join(script_dir, "data", "older"),
        os.path.join(script_dir, "..", "Data"),
        os.path.join(script_dir, "..", "data"),
    ]
    
    all_csvs = []
    for d in candidate_dirs:
        if os.path.exists(d) and os.path.isdir(d):
            try:
                for fname in os.listdir(d):
                    if fname.lower().endswith('.csv'):
                        all_csvs.append(os.path.join(d, fname))
            except Exception:
                pass

    collateral_candidates = []
    treasury_candidates = []
    fpa_candidates = []

    for fpath in all_csvs:
        fname_lower = os.path.basename(fpath).lower()
        if target_date and target_date not in fname_lower:
            continue
            
        if 'collateral' in fname_lower and 'summary' in fname_lower:
            collateral_candidates.append(fpath)
        elif 'collateral' in fname_lower:
            collateral_candidates.append(fpath)
        elif 'treasury' in fname_lower or 'cff treasury' in fname_lower:
            treasury_candidates.append(fpath)
        elif 'fp&a' in fname_lower or 'fpa' in fname_lower or 'cff fp&a' in fname_lower:
            fpa_candidates.append(fpath)

    def pick_best(candidates, preferred_name):
        if not candidates:
            return None
        # Check exact preferred name first
        for c in candidates:
            if os.path.basename(c).lower() == preferred_name.lower():
                return c
        # Otherwise sort by modification time (latest first)
        try:
            return max(candidates, key=os.path.getmtime)
        except Exception:
            return candidates[0]

    collateral_path = pick_best(collateral_candidates, "Collateral calculator 2 - Summary.csv")
    treasury_path = pick_best(treasury_candidates, "CFF Treasury.csv")
    fpa_path = pick_best(fpa_candidates, "CFF FP&A.csv")

    return collateral_path, treasury_path, fpa_path

def main():
    import sys
    target_date = None
    if len(sys.argv) > 1:
        for arg in sys.argv[1:]:
            arg_str = arg.strip()
            if re.match(r'^\d{4}-\d{2}-\d{2}$', arg_str):
                target_date = arg_str
                break

    script_dir = os.path.dirname(os.path.abspath(__file__))
    csv_dir = os.path.join(script_dir, "Data")
    html_path = os.path.join(script_dir, "collateral_dashboard.html")

    if target_date:
        print(f"Target date specified: {target_date}")

    collateral_path, treasury_path, fpa_path = find_input_csvs(script_dir, target_date)

    if not collateral_path:
        print(f"\n[NOTICE] No Collateral calculator CSV file found in candidate directories.")
        if os.path.exists(html_path) or os.path.exists(os.path.join(script_dir, "index.html")):
            print("         Existing dashboard HTML is present. Preserving current dashboard without failure.")
            # Ensure index.html exists locally
            local_index_path = os.path.join(script_dir, "index.html")
            if os.path.exists(html_path) and not os.path.exists(local_index_path):
                import shutil
                shutil.copy2(html_path, local_index_path)
            return 0
        else:
            print("ERROR: Neither CSV data nor existing HTML dashboard found.")
            return 1

    print(f"\n{'='*60}")
    print("GUARANTEE DASHBOARD UPDATE")
    print(f"{'='*60}\n")

    print(f"Parsing data:")
    collateral_data = parse_collateral_csv(collateral_path)

    treasury_data = {}
    fpa_data = {}
    if treasury_path:
        treasury_data = parse_treasury_csv(treasury_path)
    if fpa_path:
        fpa_data = parse_fpa_csv(fpa_path)

    # Bootstrap cash flow
    current_month = collateral_data['as_of_date'][:7] if collateral_data.get('as_of_date') else None
    monthly_cash, cash_sources = bootstrap_cash_flow(treasury_data, fpa_data, current_month)

    print(f"\n[OK] Data parsed successfully")
    print(f"  - Collateral as of: {collateral_data['as_of_date']}")
    print(f"  - FX Rate (EUR/NOK): {collateral_data['fx_rate']}")
    print(f"  - Utilized Today: {collateral_data['utilized_today']:,.0f} kNOK")
    print(f"  - Cash forecasts: {len(monthly_cash)} months")
    print(f"  - Cash collateral periods: {len(collateral_data['cash_collateral'])}")

    print(f"\nReading HTML...")
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    print(f"[OK] HTML read ({len(html_content):,} bytes)")

    print(f"\nUpdating HTML...")
    updated_html = update_html(html_content, collateral_data, monthly_cash, collateral_data['cash_collateral'], cash_sources)
    print(f"[OK] HTML updated")

    print(f"\nWriting updated HTML...")
    with open(html_path, 'w', encoding='utf-8') as f:
        f.write(updated_html)
    print(f"[OK] HTML written")

    # Create local index.html copy
    import shutil
    local_index_path = os.path.join(script_dir, "index.html")
    try:
        shutil.copy2(html_path, local_index_path)
        print(f"[OK] Copied collateral_dashboard.html to index.html in local folder")
    except Exception as e:
        print(f"Warning: Failed to create local index.html: {e}")

    # Maintain fixed-name copies in local Data directory
    local_data_copies = [
        (collateral_path, "Collateral calculator 2 - Summary.csv"),
        (treasury_path, "CFF Treasury.csv"),
        (fpa_path, "CFF FP&A.csv")
    ]
    for src_path, target_name in local_data_copies:
        if src_path and os.path.exists(src_path):
            dst_path = os.path.join(csv_dir, target_name)
            if os.path.abspath(src_path) != os.path.abspath(dst_path):
                try:
                    shutil.copy2(src_path, dst_path)
                    print(f"[OK] Copied {os.path.basename(src_path)} -> Data/{target_name} locally")
                except Exception as e:
                    print(f"Warning: Failed to copy {target_name} locally: {e}")

    pptx_path = os.path.join(script_dir, "collateral_dashboard.pptx")
    print(f"\nGenerating PPTX Presentation...")
    generate_pptx_dashboard(collateral_data, monthly_cash, collateral_data['cash_collateral'], pptx_path)
    print(f"[OK] PPTX written")

    # Copy files to Shared Drive if accessible
    shared_dir = u"G:\\Shared drives\\Tibber \u2013\u00a0House of Business\\Treasury\\Dashboards"
    print(f"\nChecking Shared Drive directory: {shared_dir}...")
    if os.path.exists(shared_dir):
        # 1. Main HTML dashboard remains in the root Dashboards folder
        src_html = os.path.join(script_dir, "collateral_dashboard.html")
        dst_html = os.path.join(shared_dir, "collateral_dashboard.html")
        if os.path.exists(src_html):
            try:
                shutil.copy2(src_html, dst_html)
                print(f"[OK] Copied collateral_dashboard.html to Shared Drive root")
            except Exception as e:
                print(f"Warning: Failed to copy collateral_dashboard.html to Shared Drive: {e}")

        # 2. Shared Drive Data subfolder
        shared_data_dir = os.path.join(shared_dir, "Data")
        os.makedirs(shared_data_dir, exist_ok=True)

        # Relocate collateral_dashboard.pptx and index.html to Shared Drive Data folder
        for file_name in ["collateral_dashboard.pptx", "index.html"]:
            src = os.path.join(script_dir, file_name)
            dst = os.path.join(shared_data_dir, file_name)
            if os.path.exists(src):
                try:
                    shutil.copy2(src, dst)
                    print(f"[OK] Copied {file_name} -> Data/{file_name} in Shared Drive")
                except Exception as e:
                    print(f"Warning: Failed to copy {file_name} to Shared Drive Data folder: {e}")

            # Remove old file from Shared Drive root if it exists
            old_root_file = os.path.join(shared_dir, file_name)
            if os.path.exists(old_root_file):
                try:
                    os.remove(old_root_file)
                    print(f"[OK] Relocated (removed) {file_name} from Shared Drive root")
                except Exception as e:
                    print(f"Warning: Failed to remove {file_name} from Shared Drive root: {e}")

        # 3. Copy and rename latest data files to Shared Drive Data folder
        data_copies = [
            (collateral_path, "Collateral calculator 2 - Summary.csv"),
            (treasury_path, "CFF Treasury.csv"),
            (fpa_path, "CFF FP&A.csv")
        ]
        for src_path, target_name in data_copies:
            if src_path and os.path.exists(src_path):
                dst_path = os.path.join(shared_data_dir, target_name)
                try:
                    shutil.copy2(src_path, dst_path)
                    print(f"[OK] Copied {os.path.basename(src_path)} -> Data/{target_name} in Shared Drive")
                except Exception as e:
                    print(f"Warning: Failed to copy {target_name} to Shared Drive Data folder: {e}")

        # 4. Push to GitHub if Git repository is configured
        for repo_dir in [shared_data_dir, script_dir]:
            if os.path.exists(os.path.join(repo_dir, ".git")):
                import subprocess
                try:
                    print(f"\nGit repository detected in {repo_dir}. Staging, committing, and pushing to GitHub...")
                    subprocess.run(["git", "add", "."], cwd=repo_dir, check=True)
                    status = subprocess.run(["git", "status", "--porcelain"], cwd=repo_dir, capture_output=True, text=True)
                    if status.stdout.strip():
                        commit_msg = f"Update dashboard & data: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
                        subprocess.run(["git", "commit", "-m", commit_msg], cwd=repo_dir, check=True)
                        subprocess.run(["git", "push"], cwd=repo_dir, check=True)
                        print(f"[OK] Successfully pushed changes to GitHub from {repo_dir}")
                    else:
                        print(f"[OK] No new changes to commit in {repo_dir}")
                except Exception as e:
                    print(f"Warning: Git push failed in {repo_dir}: {e}")
    else:
        print(f"Warning: Shared Drive path not found: {shared_dir}")

    print(f"\n{'='*60}")
    print(f"DASHBOARD UPDATE SUCCESSFUL")
    print(f"{'='*60}\n")
    print(f"Dashboard: {html_path}")
    print(f"Updated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")

    return 0

if __name__ == '__main__':
    try:
        exit(main())
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
        exit(1)
